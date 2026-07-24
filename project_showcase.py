"""
project_showcase.py — 產出「InsightFlow 專案實作報告」PPT。
沿用 ppt_generator 的視覺系統（色盤／卡片／徽章），輸出到 output/專案實作報告/。
一次性簡報：解決什麼問題→解決方案→架構→產出→效益→待改善→未來→下一步。
"""
import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

import ppt_generator as pg

W, H, ML, CW = pg.W, pg.H, pg.ML, pg.CW
C_RED, C_SLATE, C_SLATE2 = pg.C_RED, pg.C_SLATE, pg.C_SLATE2
C_BLUE, C_BLUE_L = pg.C_BLUE, pg.C_BLUE_L
C_INK, C_GRAY, C_WHITE, C_LIGHT, C_LINE = pg.C_INK, pg.C_GRAY, pg.C_WHITE, pg.C_LIGHT, pg.C_LINE
C_POS, C_NEG, C_NEU = pg.C_POS, pg.C_NEG, pg.C_NEU

FOOTER = "InsightFlow AI 輿情分析系統 ／ 專案實作報告"


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _title(slide, title, page, kicker=None):
    pg._rect(slide, 0, 0, Inches(0.14), Inches(1.0), fill=C_RED)
    pg._tb(slide, ML + Inches(0.15), Inches(0.22), Inches(11.3), Inches(0.62),
           title, size=27, bold=True, color=C_INK, wrap=False)
    if kicker:
        pg._tb(slide, ML + Inches(0.16), Inches(0.82), Inches(11.0), Inches(0.3),
               kicker, size=12, color=C_GRAY, wrap=False)
    pg._rect(slide, W - Inches(0.9), H - Inches(0.42), Inches(0.9), Inches(0.42), fill=C_RED)
    pg._tb(slide, W - Inches(0.9), H - Inches(0.4), Inches(0.9), Inches(0.36),
           str(page), size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, wrap=False)
    pg._tb(slide, ML, H - Inches(0.4), CW - Inches(1.2), Inches(0.32),
           FOOTER, size=9, color=C_GRAY)


def _chip(slide, x, y, w, h, title, body, accent=C_SLATE, num=None):
    """圓角資訊卡：頂部色條標題 + 內文"""
    pg._rect(slide, x, y, w, h, fill=C_LIGHT)
    pg._rect(slide, x, y, w, Inches(0.5), fill=accent)
    tx = x + Inches(0.2)
    if num is not None:
        pg._badge(slide, x + Inches(0.34), y + Inches(0.25), Inches(0.36),
                  num, fill=C_WHITE, text_color=accent, size=13)
        tx = x + Inches(0.62)
    pg._tb(slide, tx, y + Inches(0.1), w - (tx - x) - Inches(0.1), Inches(0.32),
           title, size=13, bold=True, color=C_WHITE, wrap=False)
    pg._tb(slide, x + Inches(0.2), y + Inches(0.62), w - Inches(0.4), h - Inches(0.72),
           body, size=11.5, color=C_INK)


def _arrow(slide, x, y, w, h, fill=C_SLATE2):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    shp.line.fill.background()
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    return shp


# ══════════════════════════════════════════════════════
def build(out_path):
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    # ── 1 封面 ────────────────────────────────────────
    s = _blank(prs)
    pg._diagonal(s, W - Inches(4.2), -Inches(3.8), Inches(6.0), C_LIGHT)
    pg._diagonal(s, -Inches(2.6), -Inches(2.6), Inches(4.0), C_RED)
    pg._diagonal(s, W - Inches(1.6), H - Inches(1.6), Inches(3.0), C_RED)
    pg._rect(s, Inches(0.85), Inches(1.5), Inches(2.9), Inches(0.4), fill=C_INK)
    pg._tb(s, Inches(0.85), Inches(1.56), Inches(2.9), Inches(0.3),
           "PROJECT IMPLEMENTATION", size=11, bold=True, color=C_WHITE,
           align=PP_ALIGN.CENTER, wrap=False)
    pg._tb(s, Inches(0.85), Inches(2.15), Inches(11.5), Inches(1.1),
           "InsightFlow", size=48, bold=True, color=C_INK, wrap=False)
    pg._tb(s, Inches(0.85), Inches(3.35), Inches(11.0), Inches(0.6),
           "AI 輿情分析系統 ── 專案實作報告", size=22, color=C_RED, wrap=False)
    pg._rect(s, Inches(0.85), Inches(4.2), Inches(9.5), Inches(0.03), fill=C_LINE)
    pg._tb(s, Inches(0.85), Inches(4.35), Inches(11.0), Inches(0.4),
           "一次輸入品牌、排除條件與期間，自動產出詳細輿情報告、簡報與分析架構腳本。",
           size=14, color=C_GRAY)
    pg._tb(s, Inches(0.85), Inches(4.95), Inches(10.0), Inches(0.36),
           "資料引擎｜KEYPO 大數據關鍵引擎　·　AI 敘事｜可切換後端（Ollama／Claude／OpenAI）",
           size=12, color=C_GRAY)

    # ── 2 目錄 ────────────────────────────────────────
    s = _blank(prs)
    _title(s, "報告大綱", 2)
    agenda = [
        ("解決什麼問題", "從手動多步流程到一鍵產出"),
        ("系統架構與運作", "輸入 → 抓取 → 解析 → AI → 輸出"),
        ("產出什麼成果", "報告、簡報、圖表、分析架構腳本"),
        ("帶來什麼幫助", "省時、免手動、口徑一致"),
        ("待改善與下一步", "換模型、UI 介面、分析師試用"),
    ]
    y = Inches(1.5)
    for i, (t, d) in enumerate(agenda, 1):
        ry = y + (i - 1) * Inches(1.02)
        pg._badge(s, ML + Inches(0.45), ry + Inches(0.35), Inches(0.62), i,
                  fill=C_RED if i % 2 else C_SLATE, size=18)
        pg._tb(s, ML + Inches(1.05), ry + Inches(0.06), Inches(6.5), Inches(0.4),
               t, size=17, bold=True, color=C_INK, wrap=False)
        pg._tb(s, ML + Inches(1.05), ry + Inches(0.5), Inches(9.5), Inches(0.32),
               d, size=12.5, color=C_GRAY, wrap=False)

    # ── 3 問題背景（以前的手動流程 + 痛點）────────────
    s = _blank(prs)
    _title(s, "解決什麼問題", 3, kicker="以前｜每一次分析都要手動處理，還要下載表格貼進簡報")
    steps = [
        "手寫品牌關鍵字（布林查詢語法）",
        "設定排除雜訊字串",
        "設定排除抽獎文字串",
        "上 KEYPO 設定觀測期間",
        "逐一查看聲量、情緒、關鍵字、KOL",
        "下載各項表格",
        "手動貼進簡報、重新排版",
    ]
    lx, lw = ML, Inches(7.3)
    y0 = Inches(1.45)
    for i, st in enumerate(steps, 1):
        ry = y0 + (i - 1) * Inches(0.68)
        pg._badge(s, lx + Inches(0.28), ry + Inches(0.22), Inches(0.4), i,
                  fill=C_SLATE2, size=12)
        pg._tb(s, lx + Inches(0.62), ry + Inches(0.04), lw - Inches(0.7), Inches(0.4),
               st, size=13.5, color=C_INK, wrap=False)
    # 痛點面板（右）
    px = Inches(8.2)
    pg._rect(s, px, y0, Inches(4.6), Inches(4.7), fill=C_LIGHT)
    pg._rect(s, px, y0, Inches(4.6), Inches(0.6), fill=C_NEG)
    pg._tb(s, px, y0 + Inches(0.11), Inches(4.6), Inches(0.4),
           "痛點", size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, wrap=False)
    pg._bullet_box(s, px + Inches(0.3), y0 + Inches(0.85), Inches(4.1), Inches(3.6),
                   ["每次分析都得從頭手動重來",
                    "步驟多、耗時，常要花上半天",
                    "手動操作容易漏字、出錯",
                    "各報告口徑與格式難以一致",
                    "下載表格再貼簡報，最耗工"],
                   size=13.5, spacing=12, color=C_INK)

    # ── 4 解決方案（before → after）───────────────────
    s = _blank(prs)
    _title(s, "解決方案：一次輸入，自動產出", 4,
           kicker="現在｜InsightFlow 把整條手動流程收斂成三個輸入")
    # 左：以前
    bx, bw = ML, Inches(4.4)
    pg._rect(s, bx, Inches(1.5), bw, Inches(5.1), fill=C_LIGHT)
    pg._tb(s, bx, Inches(1.62), bw, Inches(0.4), "以前", size=16, bold=True,
           color=C_GRAY, align=PP_ALIGN.CENTER, wrap=False)
    pg._bullet_box(s, bx + Inches(0.35), Inches(2.15), bw - Inches(0.6), Inches(4.2),
                   ["手寫關鍵字語法", "手設排除雜訊", "手設排除抽獎",
                    "上 KEYPO 設期間", "逐項查看數據", "下載表格", "貼簡報排版"],
                   size=13.5, spacing=9, color=C_GRAY)
    # 中：箭頭
    _arrow(s, Inches(5.15), Inches(3.55), Inches(1.1), Inches(1.0), fill=C_RED)
    # 右：現在
    rx, rw = Inches(6.55), Inches(6.28)
    pg._tb(s, rx, Inches(1.55), rw, Inches(0.4), "現在｜只要三個輸入",
           size=16, bold=True, color=C_INK, wrap=False)
    inputs = [("①", "品牌名稱", C_SLATE), ("②", "排除雜訊／抽獎", C_SLATE2),
              ("③", "時間段", C_BLUE)]
    for i, (n, lbl, acc) in enumerate(inputs):
        cy = Inches(2.15) + i * Inches(0.82)
        pg._rect(s, rx, cy, rw, Inches(0.68), fill=C_LIGHT)
        pg._badge(s, rx + Inches(0.42), cy + Inches(0.34), Inches(0.44), n,
                  fill=acc, size=14)
        pg._tb(s, rx + Inches(0.85), cy + Inches(0.14), rw - Inches(1.0), Inches(0.4),
               lbl, size=15, bold=True, color=C_INK, wrap=False)
    pg._data_card(s, rx, Inches(4.9), rw, Inches(1.55),
                  "一鍵處理", "系統自動抓取 → 分析 → 產出詳細報告",
                  accent=C_RED, value_size=26)

    # ── 5 系統架構（pipeline）─────────────────────────
    s = _blank(prs)
    _title(s, "系統架構與運作流程", 5, kicker="從輸入到成品，全自動一條龍")
    stages = [
        ("輸入", "品牌／排除／期間", C_SLATE),
        ("KEYPO 抓取", "5 端點取數", C_SLATE2),
        ("解析", "情緒／關鍵字／文章／KOL", C_SLATE),
        ("圖表", "matplotlib 產圖", C_SLATE2),
        ("AI 敘事", "可切換後端", C_BLUE),
        ("輸出", "報告／PPT／腳本", C_RED),
    ]
    n = len(stages)
    gap = Inches(0.18)
    bw2 = Emu(int((CW - gap * (n - 1)) / n))
    yb = Inches(2.4)
    for i, (t, d, acc) in enumerate(stages):
        x = ML + i * (bw2 + gap)
        pg._rect(s, x, yb, bw2, Inches(1.5), fill=acc)
        pg._tb(s, x + Inches(0.05), yb + Inches(0.34), bw2 - Inches(0.1), Inches(0.5),
               t, size=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, wrap=True)
        pg._tb(s, x + Inches(0.05), yb + Inches(0.92), bw2 - Inches(0.1), Inches(0.5),
               d, size=9.5, color=C_WHITE, align=PP_ALIGN.CENTER, wrap=True)
        if i < n - 1:
            pg._tb(s, x + bw2 - Inches(0.02), yb + Inches(0.45), gap + Inches(0.1),
                   Inches(0.5), "›", size=20, bold=True, color=C_GRAY,
                   align=PP_ALIGN.CENTER, wrap=False)
    pg._tb(s, ML, Inches(4.4), CW, Inches(0.35),
           "▎ 對應模組", size=13, bold=True, color=C_INK)
    pg._bullet_box(s, ML, Inches(4.85), CW, Inches(1.9),
                   ["keypo_fetcher（呼叫 KEYPO CLI）　·　keypo_loader（解析）　·　"
                    "chart_generator（圖表）",
                    "ai_backend（後端切換）＋ local_ai_analyzer（六節分析、事件敘事）",
                    "detailed_report_generator（詳細報告）　·　ppt_generator（簡報）　·　"
                    "framework_generator（分析架構腳本）"],
                   size=12.5, spacing=8, color=C_GRAY)

    # ── 6 產出成果（4 卡）─────────────────────────────
    s = _blank(prs)
    _title(s, "一次執行，產出四種成果", 6, kicker="全部自動存進 output／{品牌}／ 資料夾")
    cards = [
        ("詳細分析報告", "§1–§12 完整 Markdown：執行摘要、聲量趨勢（含事件敘事）、"
         "情緒、關鍵字、熱門文章、KOL、風險、競品。", C_SLATE),
        ("簡報 PPT", "26 頁 KEYPO 風格：趨勢事件、來源分布、好感度、關鍵字對應事件、"
         "KOL 四維、競品分析、觀測彙整。", C_RED),
        ("視覺圖表", "聲量趨勢（事件標註）、情緒甜甜圈、來源分布、關鍵字長條與文字雲、"
         "競品比較，皆自動嵌入簡報。", C_SLATE2),
        ("分析架構與腳本", "Word：AI 依品牌自動擬定議題×維度矩陣、KEYPO 對應、"
         "查詢式原則、逐議題腳本、章節對應。", C_BLUE),
    ]
    gx = Inches(0.25)
    cw = Emu(int((CW - gx) / 2))
    ch = Inches(2.35)
    for i, (t, b, acc) in enumerate(cards):
        col, row = i % 2, i // 2
        x = ML + col * (cw + gx)
        y = Inches(1.55) + row * (ch + Inches(0.25))
        _chip(s, x, y, cw, ch, t, b, accent=acc, num=i + 1)

    # ── 7 效益 ────────────────────────────────────────
    s = _blank(prs)
    _title(s, "帶來什麼幫助", 7)
    pg._kpi_row(s, ML, Inches(1.55), CW, Inches(1.35), [
        ("半天 → 分鐘", "產製時間", C_SLATE),
        ("0", "手動複製貼上", C_RED),
        ("1 鍵", "完成整條流程", C_SLATE2),
        ("100%", "口徑與格式一致", C_BLUE),
    ], value_size=22)
    rows = [
        ("省時", "手動半天的取數、製表、排版，收斂成一次執行、數分鐘完成。"),
        ("免手動複製貼上", "圖表與表格自動產生並嵌入簡報，不再下載 → 貼上 → 重排。"),
        ("口徑一致", "抽獎排除、雜訊排除、期間比較（MoM／YoY）每份報告統一。"),
        ("AI 事件敘事", "聲量高峰自動對應具體事件，並產出洞察與建議，不只是數字。"),
        ("可切換 AI 後端", "Ollama（免費離線）／Claude／OpenAI，改一個設定即可切換。"),
    ]
    y = Inches(3.25)
    for i, (h, d) in enumerate(rows):
        ry = y + i * Inches(0.72)
        pg._oval(s, ML, ry + Inches(0.04), Inches(0.42), Inches(0.42),
                 fill=(C_RED if i % 2 == 0 else C_SLATE))
        pg._tb(s, ML + Inches(0.6), ry - Inches(0.02), Inches(2.6), Inches(0.4),
               h, size=14.5, bold=True, color=C_INK, wrap=False)
        pg._tb(s, ML + Inches(3.1), ry + Inches(0.02), Inches(9.1), Inches(0.6),
               d, size=12.5, color=C_GRAY, wrap=True)

    # ── 8 待改善（換模型）─────────────────────────────
    s = _blank(prs)
    _title(s, "待改善：AI 模型", 8, kicker="已建可切換後端架構，換模型只需改一個設定")
    # 現況
    ax, aw = ML, Inches(5.9)
    pg._rect(s, ax, Inches(1.55), aw, Inches(4.9), fill=C_LIGHT)
    pg._rect(s, ax, Inches(1.55), aw, Inches(0.55), fill=C_NEU)
    pg._tb(s, ax, Inches(1.65), aw, Inches(0.4), "現況：Ollama（本地免費）",
           size=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, wrap=False)
    pg._bullet_box(s, ax + Inches(0.35), Inches(2.3), aw - Inches(0.6), Inches(4.0),
                   ["模型：qwen2.5:7b，離線、零成本、資料留本機",
                    "限制①：偶爾漏出簡體字（已加 OpenCC 簡→繁保險）",
                    "限制②：內容深度有限、偶有小幻覺",
                    "適合：測試、離線、對隱私敏感的情境"],
                   size=13, spacing=12, color=C_INK)
    # 升級
    bx2 = Inches(6.9)
    pg._rect(s, bx2, Inches(1.55), aw, Inches(4.9), fill=C_BLUE_L)
    pg._rect(s, bx2, Inches(1.55), aw, Inches(0.55), fill=C_BLUE)
    pg._tb(s, bx2, Inches(1.65), aw, Inches(0.4), "升級：雲端模型",
           size=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, wrap=False)
    pg._bullet_box(s, bx2 + Inches(0.35), Inches(2.3), aw - Inches(0.6), Inches(4.0),
                   ["換 GPT-4o / Claude：繁中更乾淨、議題與敘事更準更細",
                    "程式已就緒：改 .env 的 AI_BACKEND 即可切換",
                    "成本極低：每份報告約數美分",
                    "待決策：依成本／隱私／品質選定正式後端"],
                   size=13, spacing=12, color=C_INK)

    # ── 9 未來提升 ────────────────────────────────────
    s = _blank(prs)
    _title(s, "未來系統提升方向", 9, kicker="從單品牌報告器，邁向季報等級分析引擎")
    grid = [
        ("資料端點升級補強", "來源結構、KOL 現已有（推估版）；升級為官方 datadist"
         "（全站聲量）、opleader、hotkwV5，並新增 hotchnl 頻道／通路排行，精度更高。"),
        ("多議題 × 多品牌", "以「議題 × 維度」矩陣統一結構，支援多品牌競爭格局，"
         "對標專業季報。"),
        ("期間比較升級", "QoQ 季比較、近 12 個月趨勢線，補足現有 MoM／YoY。"),
        ("架構腳本資料驅動", "分析架構腳本從『AI 規劃』升級為『真實資料回填』，"
         "直接產出可交付的分析草稿。"),
    ]
    gx = Inches(0.25)
    cw = Emu(int((CW - gx) / 2))
    ch = Inches(2.3)
    for i, (t, b) in enumerate(grid):
        col, row = i % 2, i // 2
        x = ML + col * (cw + gx)
        y = Inches(1.55) + row * (ch + Inches(0.25))
        acc = [C_SLATE, C_RED, C_SLATE2, C_BLUE][i]
        _chip(s, x, y, cw, ch, t, b, accent=acc, num=i + 1)

    # ── 10 下一步 ─────────────────────────────────────
    s = _blank(prs)
    _title(s, "下一步", 10, kicker="讓系統更好用、更貼近分析師實際需求")
    nexts = [
        ("開發 UI 介面的直覺系統",
         "以圖形介面取代命令列輸入，讓非工程背景的分析師也能直覺操作："
         "填欄位、按執行、看結果、下載成品。降低使用門檻。", C_RED),
        ("交付分析師試用並蒐集回饋",
         "請分析師實際使用，蒐集需求與需修改處（版型、口徑、欄位、"
         "議題邏輯等），依回饋迭代，逐步貼近正式提報標準。", C_SLATE),
    ]
    y = Inches(1.7)
    for i, (t, b, acc) in enumerate(nexts):
        ry = y + i * Inches(2.35)
        pg._rect(s, ML, ry, CW, Inches(2.05), fill=C_LIGHT)
        pg._badge(s, ML + Inches(0.7), ry + Inches(1.0), Inches(0.9), i + 1,
                  fill=acc, size=26)
        pg._tb(s, ML + Inches(1.5), ry + Inches(0.28), Inches(10.3), Inches(0.5),
               t, size=19, bold=True, color=C_INK, wrap=False)
        pg._tb(s, ML + Inches(1.5), ry + Inches(0.9), Inches(10.3), Inches(1.0),
               b, size=13.5, color=C_GRAY, wrap=True)

    # ── 11 封底 ───────────────────────────────────────
    s = _blank(prs)
    pg._diagonal(s, -Inches(1.8), -Inches(3.8), Inches(6.0), C_LIGHT)
    pg._diagonal(s, W - Inches(1.4), H - Inches(1.4), Inches(4.0), C_RED)
    pg._diagonal(s, -Inches(2.4), H - Inches(1.6), Inches(3.0), C_RED)
    pg._tb(s, Inches(0.9), Inches(2.6), Inches(11.0), Inches(1.1),
           "InsightFlow", size=44, bold=True, color=C_INK, wrap=False)
    pg._tb(s, Inches(0.9), Inches(3.75), Inches(11.0), Inches(0.6),
           "讓輿情分析，從半天手工變成一鍵產出。", size=20, color=C_RED, wrap=False)
    pg._rect(s, Inches(0.9), Inches(4.5), Inches(8.5), Inches(0.03), fill=C_LINE)
    pg._tb(s, Inches(0.9), Inches(4.65), Inches(10.0), Inches(0.4),
           "AI 輿情分析系統 ／ 專案實作報告", size=13, color=C_GRAY)
    pg._tb(s, Inches(0.9), Inches(6.5), Inches(10.0), Inches(0.35),
           "— 簡報完 —", size=12, color=C_SLATE)

    os.makedirs(os.path.dirname(out_path) or ".", exist_ok=True)
    prs.save(out_path)
    return out_path


if __name__ == "__main__":
    out = build("output/專案實作報告/InsightFlow_專案實作報告.pptx")
    print(f"已產出：{out}")
