"""
ppt_generator.py — InsightFlow PPT 產生器（KEYPO 官方報告風格）
12 頁結構：封面 / 分節頁 / 執行摘要 / 聲量趨勢(callout圖) /
           情緒分布(甜甜圈) / 熱門關鍵字 / 關鍵字雲 /
           熱門文章(TOP3卡片+表格) / KOL(排名徽章) / AI洞察 /
           結論(3欄卡) / 封底
           ── 若提供 competitor_data，KOL 頁後會插入「競品聲量分析」章節
           （分節頁 + 競品聲量比較 + 競品情緒比較 + 競品熱門文章，共 4 頁）
"""
import os
import re
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from keypo_loader import (
    classify_article,
    parse_kol_from_articles,
    match_keywords_to_articles,
)


# ── 色盤（KEYPO 官方報告風格）───────────────────────────
C_RED    = RGBColor(0xCC, 0x32, 0x32)   # 紅   — 強調色、對角裝飾、頁碼色塊
C_SLATE  = RGBColor(0x52, 0x67, 0x83)   # 深藍灰 — KPI 卡、表格標頭
C_SLATE2 = RGBColor(0x82, 0x95, 0xAF)   # 藍灰 — 表格標頭（次要）、排名徽章
C_BLUE   = RGBColor(0x43, 0x71, 0xC3)   # 藍   — 結論卡標題列
C_BLUE_L = RGBColor(0xD1, 0xD6, 0xEC)   # 淺藍 — 結論卡卡身
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_INK    = RGBColor(0x1A, 0x1A, 0x1A)   # 內文深色字
C_GRAY   = RGBColor(0x6B, 0x6B, 0x6B)   # 次要文字灰
C_LIGHT  = RGBColor(0xF4, 0xF6, 0xF8)   # 卡片/表格底色
C_LINE   = RGBColor(0xD5, 0xD9, 0xDF)   # 分隔線

# 情緒語意色（沿用專案既有色彩慣例，不隨版型更動）
C_POS    = RGBColor(0x3B, 0x7D, 0xD8)   # 藍 — 正面
C_NEG    = RGBColor(0xC0, 0x39, 0x2B)   # 紅 — 負面
C_NEU    = RGBColor(0x7A, 0x7D, 0x82)   # 灰 — 中立

# 文章分類標籤色（新聞媒體／討論區／官方／網友分享）
CATEGORY_COLORS = {
    "新聞媒體": RGBColor(0x52, 0x67, 0x83),
    "討論區":   RGBColor(0x43, 0x71, 0xC3),
    "官方":     RGBColor(0xCC, 0x32, 0x32),
    "網友分享": RGBColor(0x82, 0x95, 0xAF),
}

FONT = "Microsoft JhengHei"

# ── 版面常數 ─────────────────────────────────────────
W      = Inches(13.33)   # 投影片寬
H      = Inches(7.5)     # 投影片高
ML     = Inches(0.5)     # 左邊距
CW     = Inches(12.33)   # 內容寬 (W - 2*ML)
CY     = Inches(1.0)     # 內容起始 Y（無滿版頁首後，標題正下方）

CIRCLED_DIGITS = ["①", "②", "③", "④", "⑤", "⑥", "⑦", "⑧", "⑨", "⑩"]


# ══════════════════════════════════════════════════════
# 基礎工具
# ══════════════════════════════════════════════════════

def _rect(slide, x, y, w, h, fill=None):
    """加入矩形色塊"""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.line.fill.background()
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    return shp


def _oval(slide, x, y, w, h, fill=None, line=False):
    """加入橢圓/圓形色塊"""
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    if line:
        shp.line.color.rgb = fill if fill else C_LINE
        shp.line.width = Pt(1.5)
        shp.fill.background()
    else:
        shp.line.fill.background()
        if fill:
            shp.fill.solid()
            shp.fill.fore_color.rgb = fill
        else:
            shp.fill.background()
    return shp


def _diagonal(slide, x, y, size, fill):
    """加入 45 度旋轉的正方形色塊，用於對角裝飾"""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, size, size)
    shp.rotation = 45
    shp.line.fill.background()
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    return shp


def _tb(slide, x, y, w, h, text,
        size=13, bold=False, color=C_INK,
        align=PP_ALIGN.LEFT, wrap=True):
    """加入單段文字框"""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = str(text)
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def _bullet_box(slide, x, y, w, h, items,
                size=13, spacing=6, color=C_INK):
    """加入項目符號清單"""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        if not first:
            p.space_before = Pt(spacing)
        first = False
        r = p.add_run()
        r.text = str(item)
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.color.rgb = color


def _numbered_insight(slide, x, y, w, h, items, size=13, spacing=8, color=C_INK):
    """圈碼列點洞察（①②③…，無底色框，直接鋪在標題下方，仿範本樣式）"""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        if not first:
            p.space_before = Pt(spacing)
        first = False
        digit = CIRCLED_DIGITS[i] if i < len(CIRCLED_DIGITS) else f"{i + 1}."
        r = p.add_run()
        r.text = f"{digit} {item}"
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return box


def _cell(cell, text, size=11, bold=False,
          color=C_INK, bg=None, align=PP_ALIGN.LEFT):
    """設定表格儲存格"""
    if bg:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
    cell.text = str(text)
    tf = cell.text_frame
    tf.word_wrap = True
    for para in tf.paragraphs:
        para.alignment = align
        for run in para.runs:
            run.font.name = FONT
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color


def _clean(text, max_len=40):
    """移除 emoji 與換行，截斷過長文字"""
    cleaned = ''.join(c for c in str(text) if ord(c) <= 0xFFFF)
    cleaned = re.sub(r'[\x00-\x1f]', ' ', cleaned)
    cleaned = ' '.join(cleaned.split())
    if len(cleaned) > max_len:
        cleaned = cleaned[:max_len] + '…'
    return cleaned


def _parse_ai(text):
    """解析 AI 報告各節（# 標題）"""
    sections = {}
    key, lines = None, []
    for line in text.split('\n'):
        if line.startswith('# '):
            if key:
                sections[key] = '\n'.join(lines).strip()
            key, lines = line[2:].strip(), []
        else:
            lines.append(line)
    if key:
        sections[key] = '\n'.join(lines).strip()
    return sections


def _strip_md(text):
    """移除 Markdown 粗體、清單符號"""
    text = re.sub(r'\*\*([^*]+)\*\*', r'\1', text)
    text = re.sub(r'\*([^*]+)\*', r'\1', text)
    return text


def format_date(s):
    return f"{s[:4]}/{s[4:6]}/{s[6:]}"


def _safe_pct(s):
    """把 '32.53%' 這類百分比字串轉為 float（32.53）；失敗回 0.0"""
    try:
        return float(str(s).replace('%', '').strip())
    except (TypeError, ValueError):
        return 0.0


# ══════════════════════════════════════════════════════
# 共用版面元件
# ══════════════════════════════════════════════════════

def _header(slide, title, page_num, brand, period):
    """白底頁首：左側紅色色塊 + 標題 + 右下角紅色頁碼色塊 + 頁尾"""
    _rect(slide, 0, 0, Inches(0.14), Inches(1.0), fill=C_RED)
    _tb(slide, ML + Inches(0.15), Inches(0.2),
        Inches(11.0), Inches(0.6),
        title, size=26, bold=True, color=C_INK, wrap=False)

    _rect(slide, W - Inches(0.9), H - Inches(0.42), Inches(0.9), Inches(0.42), fill=C_RED)
    _tb(slide, W - Inches(0.9), H - Inches(0.4),
        Inches(0.9), Inches(0.36),
        str(page_num), size=12, bold=True, color=C_WHITE,
        align=PP_ALIGN.CENTER, wrap=False)

    _tb(slide, ML, H - Inches(0.4), CW - Inches(1.2), Inches(0.32),
        f"{brand} 網路輿情分析報告  ／  {period}  ／  資料來源：KEYPO 大數據關鍵引擎",
        size=9, color=C_GRAY)


def _data_card(slide, x, y, w, h, value, label, accent=C_SLATE, value_size=24):
    """KPI 卡片（純色塊：標籤 + 大數字，白字置中）"""
    _rect(slide, x, y, w, h, fill=accent)
    _tb(slide, x + Inches(0.08), y + Inches(0.1),
        w - Inches(0.16), Inches(0.3),
        label, size=10, bold=True, color=C_WHITE,
        align=PP_ALIGN.CENTER, wrap=False)
    _tb(slide, x + Inches(0.08), y + Inches(0.42),
        w - Inches(0.16), h - Inches(0.5),
        value, size=value_size, bold=True, color=C_WHITE,
        align=PP_ALIGN.CENTER, wrap=False)


def _kpi_row(slide, x, y, w, h, items, value_size=19):
    """一列等寬 KPI 卡（items: [(value, label, accent), ...]）"""
    n = len(items)
    gap = Inches(0.15)
    card_w = Emu(int((w - gap * (n - 1)) / n))
    for i, (value, label, accent) in enumerate(items):
        _data_card(
            slide,
            x + i * (card_w + gap), y,
            card_w, h, value, label,
            accent=accent, value_size=value_size
        )


RANK_COLORS = [C_RED, C_SLATE, C_SLATE2, C_NEU]


def _badge(slide, cx, cy, d, number, fill=C_SLATE, text_color=C_WHITE, size=13):
    """圓形排名徽章（置中數字）"""
    x, y = cx - d // 2, cy - d // 2
    _oval(slide, x, y, d, d, fill=fill)
    box = slide.shapes.add_textbox(x, y, d, d)
    tf = box.text_frame
    tf.word_wrap = False
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = str(number)
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = text_color
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return box


def _rank_badge(slide, cx, cy, d, rank, size=13):
    """依名次取色的排名徽章（前三名用強調色，其餘中性灰）"""
    color = RANK_COLORS[rank - 1] if rank <= 3 else RANK_COLORS[3]
    _badge(slide, cx, cy, d, rank, fill=color, size=size)


def _section_divider(prs, title, subtitle, chapter, page_num=None):
    """大字報式分節頁（白底 + 對角裝飾 + CHAPTER 字樣）。
    chapter 為章節序號（CHAPTER N），page_num 為連續頁碼（右下徽章）。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _diagonal(slide, W - Inches(3.6), -Inches(3.3), Inches(5.4), C_LIGHT)
    _diagonal(slide, -Inches(2.2), H - Inches(2.0), Inches(3.4), C_RED)

    _tb(slide, Inches(0.9), Inches(1.35), Inches(6.0), Inches(0.45),
        f"CHAPTER {chapter}", size=15, bold=True, color=C_SLATE2, wrap=False)
    _tb(slide, Inches(0.9), Inches(2.7), Inches(10.5), Inches(1.1),
        title, size=40, bold=True, color=C_INK)
    _rect(slide, Inches(0.9), Inches(3.55), Inches(0.9), Inches(0.06), fill=C_RED)
    _tb(slide, Inches(0.9), Inches(3.72), Inches(10.5), Inches(0.5),
        subtitle, size=14, color=C_GRAY)

    badge_num = page_num if page_num is not None else chapter
    _badge(slide, W - Inches(0.75), H - Inches(0.75), Inches(0.6),
           f"{badge_num:02d}", fill=C_RED, size=13)

    return slide


def _section_label(slide, x, y, text, color=C_INK):
    """小節標籤（▎ 字首）"""
    _tb(slide, x, y, Inches(6.0), Inches(0.32),
        f"▎ {text}", size=12, bold=True, color=color)


# ══════════════════════════════════════════════════════
# 各頁生成
# ══════════════════════════════════════════════════════

def _slide_cover(prs, brand, period):
    """封面：白底 + 紅色對角裝飾"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _diagonal(slide, W - Inches(4.2), -Inches(3.8), Inches(6.0), C_LIGHT)
    _diagonal(slide, -Inches(2.6), -Inches(2.6), Inches(4.0), C_RED)
    _diagonal(slide, W - Inches(1.6), H - Inches(1.6), Inches(3.0), C_RED)

    _rect(slide, Inches(0.85), Inches(1.55), Inches(2.6), Inches(0.38), fill=C_INK)
    _tb(slide, Inches(0.85), Inches(1.61), Inches(2.6), Inches(0.3),
        "SOCIAL LISTENING REPORT", size=11, bold=True,
        color=C_WHITE, align=PP_ALIGN.CENTER, wrap=False)

    _tb(slide, Inches(0.85), Inches(2.15), Inches(11.5), Inches(1.15),
        brand, size=42, bold=True, color=C_INK)

    _tb(slide, Inches(0.85), Inches(3.4), Inches(10.0), Inches(0.65),
        "網路輿情分析報告", size=22, color=C_RED)

    _rect(slide, Inches(0.85), Inches(4.2), Inches(9.5), Inches(0.03), fill=C_LINE)

    _tb(slide, Inches(0.85), Inches(4.35), Inches(10.0), Inches(0.45),
        f"分析期間：{period}", size=14, color=C_GRAY)

    _tb(slide, Inches(0.85), Inches(4.9), Inches(10.0), Inches(0.38),
        "資料來源｜KEYPO 大數據關鍵引擎", size=12, color=C_GRAY)


def _info_card(slide, x, y, w, h, title, body, accent=C_SLATE):
    """左側色條 + 標題 + 內文的資訊卡（供調查方法／名詞說明頁使用）"""
    _rect(slide, x, y, w, h, fill=C_LIGHT)
    _rect(slide, x, y, Inches(0.08), h, fill=accent)
    _tb(slide, x + Inches(0.24), y + Inches(0.14), w - Inches(0.4), Inches(0.34),
        title, size=13, bold=True, color=C_INK, wrap=False)
    _tb(slide, x + Inches(0.24), y + Inches(0.52), w - Inches(0.42), h - Inches(0.6),
        body, size=11, color=C_GRAY)


def _slide_method(prs, brand, period, page_num):
    """調查方法：觀測區間 / 資料來源 / 分析工具（對照官方月報開場）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, "調查方法", page_num, brand, period)

    cards = [
        ("觀測區間", period + "（依 KEYPO 大數據關鍵引擎擷取之網路討論）", C_SLATE),
        ("資料來源",
         "KEYPO 大數據關鍵引擎——彙整社群、新聞、討論區、部落格等公開網路"
         "討論文章與新聞，計算網路聲量。", C_RED),
        ("分析工具",
         "透過斷字切詞（Word Segmentation）、語意分析（Semantic Analysis）"
         "與即時情緒分析（Real-Time Sentiment Analysis）等技術，將文章"
         "分類為正面、負面、中立三種情緒。", C_BLUE),
        ("AI 輔助洞察",
         "本報告之敘事摘要與洞察由本地 AI 模型（Ollama / qwen2.5:7b）"
         "依原始資料輔助生成，數值以 KEYPO 原始資料為準。", C_SLATE2),
    ]
    y = Inches(1.15)
    for i, (t, b, acc) in enumerate(cards):
        _info_card(slide, ML, y + i * Inches(1.45), CW, Inches(1.28), t, b, accent=acc)


def _slide_glossary(prs, brand, period, page_num):
    """名詞說明：網路聲量 / 好感度 / 頻道 / 社群 / 討論區 / KOL"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, "名詞說明", page_num, brand, period)

    terms = [
        ("網路聲量", "網路討論或新聞的則數，包含主文與留言；主文標題若提及分析"
                     "對象，其留言亦計入；若僅留言提及，則僅計提及之留言則數。"),
        ("好感度（P/N）", "正負評比＝正面聲量 ÷ 負面聲量，代表正負面文章比例；"
                          "數值越高代表口碑越正向。"),
        ("頻道", "網站平台細分之討論區、社群、部落格、新聞主題分類等。"),
        ("社群", "Facebook、Instagram、YouTube、Threads、小紅書等社群平台。"),
        ("討論區", "PTT、Dcard、Mobile01 等論壇。"),
        ("關鍵領袖（KOL）", "網路議題中特別具影響力、能引領網友討論風向的人或組織。"),
    ]
    gap = Inches(0.25)
    col_w = Emu(int((CW - gap) / 2))
    card_h = Inches(1.7)
    for i, (t, b) in enumerate(terms):
        col = i % 2
        row = i // 2
        x = ML + col * (col_w + gap)
        y = Inches(1.15) + row * (card_h + Inches(0.18))
        _info_card(slide, x, y, col_w, card_h, t, b, accent=RANK_COLORS[i % 4])


def _slide_summary(prs, brand, period, trend, sentiment, keywords,
                    growth_info=None, page_num=2):
    """執行摘要：KPI 卡片 + 上期比較 + 重點列表"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, "執行摘要", page_num, brand, period)

    # ── 資料卡片 (4 個) ──────────────────────────────
    card_w = Inches(2.88)
    card_h = Inches(1.28)
    gap    = Inches(0.19)
    cy     = Inches(1.02)

    cards = [
        (f"{trend['total']:,}筆",          "本期總聲量",  C_SLATE),
        (str(sentiment["pn_value"]),        "P/N 值",     C_BLUE),
        (sentiment["positive_percent"],     "正面聲量",    C_POS),
        (sentiment["negative_percent"],     "負面聲量",    C_NEG),
    ]
    for i, (val, lbl, acc) in enumerate(cards):
        _data_card(slide,
                   ML + i * (card_w + gap), cy,
                   card_w, card_h, val, lbl, accent=acc)

    # ── 上期比較／去年同期比較（簡化版，僅在有快照時顯示）──
    comparison_lines = []
    if growth_info:
        mom = growth_info.get("mom")
        yoy = growth_info.get("yoy")
        if mom:
            sign = "+" if mom["rate"] >= 0 else ""
            comparison_lines.append(
                f"上期：{mom['prev_total']:,}筆｜本期聲量成長率 {sign}{mom['rate']}%"
            )
        if yoy:
            sign = "+" if yoy["rate"] >= 0 else ""
            comparison_lines.append(
                f"去年同期：{yoy['prev_total']:,}筆｜年增率 {sign}{yoy['rate']}%"
            )
    for i, line in enumerate(comparison_lines):
        _tb(slide, ML, cy + card_h + Inches(0.04) + i * Inches(0.2),
            card_w, Inches(0.2), line, size=9.5, color=C_GRAY, wrap=True)

    # ── 分隔 ─────────────────────────────────────────
    _tb(slide, ML, Inches(2.86), Inches(4.0), Inches(0.3),
        "▎ 本期重點摘要", size=12, bold=True, color=C_INK)
    _rect(slide, ML, Inches(3.2), CW, Inches(0.02), fill=C_LINE)

    # ── 重點列表（圈碼）──────────────────────────────
    peak = trend["peak_date"]
    neg  = sentiment["negative"]
    kw0, kw1, kw2 = keywords[0], keywords[1], keywords[2]

    bullets = [
        f"聲量高峰 {peak}（{trend['peak_volume']:,}筆），"
        f"本期累積 {trend['total']:,}筆，平均每日"
        f" {trend['total'] // len(trend['daily']):,}筆。",

        f"情緒結構：正面 {sentiment['positive_percent']}、"
        f"中立 {sentiment['neutral_percent']}、"
        f"負面 {sentiment['negative_percent']}；"
        f"P/N 值 {sentiment['pn_value']}。",

        f"熱門關鍵字 TOP3：{kw0['keyword']}（{kw0['count']}）、"
        f"{kw1['keyword']}（{kw1['count']}）、"
        f"{kw2['keyword']}（{kw2['count']}）。",

        f"負面聲量 {neg:,}筆（{sentiment['negative_percent']}），"
        + ("超過 100 筆，建議持續監控社群討論。"
           if neg > 100 else "維持低水位。"),
    ]

    _numbered_insight(slide, ML, Inches(3.36), CW, Inches(3.55),
                       bullets, size=14, spacing=14)


def _slide_trend(prs, brand, period, trend, articles, charts_dir,
                  growth_info=None, page_num=3, narrative=None):
    """聲量趨勢圖：事件驅動敘事（總聲量＋來源結構＋各高峰對應事件）＋折線圖。
    未提供 narrative 時退回模板洞察＋KPI 列。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, "聲量趨勢分析", page_num, brand, period)

    days = len(trend["daily"])

    if narrative:
        # 事件驅動敘事（第 1 點來源結構、之後各高峰事件），給予較大文字區
        _numbered_insight(slide, ML, Inches(0.95), CW, Inches(2.5),
                          narrative, size=13, spacing=11)
        slide.shapes.add_picture(
            f"{charts_dir}/volume_trend.png",
            ML, Inches(3.55), width=CW, height=Inches(3.4))
        return

    # ── 退化版：模板洞察 + KPI 列 + 折線圖 ──────────────
    total_engagement = sum(a.get("engagement", 0) for a in articles)
    insight_items = [
        f"本期{brand}總聲量共{trend['total']:,}筆，聲量高峰出現在"
        f"{trend['peak_date']}（{trend['peak_volume']:,}筆）。",
        f"日均聲量約{trend['total'] // days:,}筆。" if days else "日均聲量無法計算。",
    ]
    if growth_info:
        comparison_parts = []
        mom = growth_info.get("mom")
        yoy = growth_info.get("yoy")
        if mom:
            word = "成長" if mom["rate"] >= 0 else "下滑"
            comparison_parts.append(
                f"較上期（{mom['prev_total']:,}筆）{word} {abs(mom['rate'])}%"
            )
        if yoy:
            word = "成長" if yoy["rate"] >= 0 else "下滑"
            comparison_parts.append(
                f"較去年同期（{yoy['prev_total']:,}筆）{word} {abs(yoy['rate'])}%"
            )
        if comparison_parts:
            insight_items.append("本期聲量" + "；".join(comparison_parts) + "。")
    _numbered_insight(slide, ML, Inches(0.95), CW, Inches(1.05), insight_items, size=12)

    _kpi_row(slide, ML, Inches(2.05), CW, Inches(0.85), [
        (f"{trend['total']:,}", "總聲量（筆）", C_SLATE),
        (f"{trend['peak_volume']:,}", "聲量高峰（筆）", C_RED),
        (f"{trend['total'] // days:,}" if days else "0", "日均聲量（筆）", C_SLATE2),
        (f"{len(articles):,}", "熱門文章數", C_SLATE),
        (f"{total_engagement:,}", "累計互動數", C_SLATE2),
    ], value_size=18)

    slide.shapes.add_picture(
        f"{charts_dir}/volume_trend.png",
        ML, Inches(3.05), width=CW, height=Inches(3.9))


def _slide_sentiment(prs, brand, period, sentiment, charts_dir, page_num=4):
    """情緒分布：圈碼洞察 + 甜甜圈圖（中心 P/N 值）+ 數據列"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, "情緒分布分析", page_num, brand, period)

    dominant = max(
        [("正面", sentiment["positive"]), ("中立", sentiment["neutral"]),
         ("負面", sentiment["negative"])],
        key=lambda x: x[1]
    )[0]
    _numbered_insight(slide, ML, Inches(0.95), CW, Inches(0.6), [
        f"本期輿情以{dominant}聲量為主，P/N 值為 {sentiment['pn_value']}。"
    ], size=12)

    # 甜甜圈圖（左，中心已內嵌 P/N 值）
    slide.shapes.add_picture(
        f"{charts_dir}/sentiment_pie.png",
        ML, Inches(1.65), width=Inches(6.3), height=Inches(5.1))

    # 數據區（右）
    rx = Inches(7.1)

    _data_card(slide, rx, Inches(1.65), Inches(5.7), Inches(0.9),
               str(sentiment["pn_value"]), "P/N 值（正面聲量 ÷ 負面聲量）",
               accent=C_SLATE, value_size=30)

    rows = [
        ("正面", sentiment["positive"], sentiment["positive_percent"], C_POS),
        ("中立", sentiment["neutral"],  sentiment["neutral_percent"],  C_NEU),
        ("負面", sentiment["negative"], sentiment["negative_percent"], C_NEG),
    ]
    for i, (label, vol, pct, color) in enumerate(rows):
        sy = Inches(2.75) + i * Inches(1.25)
        _rect(slide, rx, sy + Inches(0.14), Inches(0.22), Inches(0.22), fill=color)
        _tb(slide, rx + Inches(0.38), sy,
            Inches(1.8), Inches(0.32), label, size=13, bold=True, color=C_INK)
        _tb(slide, rx + Inches(0.38), sy + Inches(0.35),
            Inches(5.0), Inches(0.46),
            f"{vol:,} 筆", size=20, bold=True, color=color, wrap=False)
        _tb(slide, rx + Inches(0.38), sy + Inches(0.84),
            Inches(5.0), Inches(0.3), pct, size=12, color=C_GRAY)


def _slide_source(prs, brand, period, source_dist, charts_dir, page_num):
    """來源分布：來源類型甜甜圈圖 + 數據列（社群／新聞／討論區／部落格）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, "來源分布分析", page_num, brand, period)

    items = source_dist["items"]
    top = max(items, key=lambda x: x["count"]) if items else None
    insight = []
    if top and top["count"]:
        insight.append(
            f"本期熱門討論主要來自「{top['type']}」（{top['percent']}%），"
            f"共 {source_dist['total']} 篇熱門文章。"
        )
    insight.append("（來源分布依熱門文章來源類型統計，為討論結構參考值。）")
    _numbered_insight(slide, ML, Inches(0.95), CW, Inches(0.9), insight, size=12)

    slide.shapes.add_picture(
        f"{charts_dir}/source_pie.png",
        ML, Inches(1.75), width=Inches(6.3), height=Inches(5.0))

    rx = Inches(7.1)
    _data_card(slide, rx, Inches(1.75), Inches(5.7), Inches(0.9),
               f"{source_dist['total']} 篇", "本期熱門文章總數",
               accent=C_SLATE, value_size=26)

    for i, it in enumerate(items):
        sy = Inches(2.9) + i * Inches(0.95)
        color = {
            "社群": C_POS, "新聞": C_SLATE,
            "討論區": C_BLUE, "部落格": C_NEU,
        }.get(it["type"], C_SLATE2)
        _rect(slide, rx, sy + Inches(0.1), Inches(0.22), Inches(0.22), fill=color)
        _tb(slide, rx + Inches(0.38), sy, Inches(2.2), Inches(0.34),
            it["type"], size=14, bold=True, color=C_INK, wrap=False)
        _tb(slide, rx + Inches(2.2), sy - Inches(0.02),
            Inches(1.7), Inches(0.4),
            f"{it['count']} 篇", size=16, bold=True, color=color, wrap=False)
        _tb(slide, rx + Inches(4.0), sy + Inches(0.02),
            Inches(1.6), Inches(0.34),
            f"{it['percent']}%", size=14, color=C_GRAY, wrap=False)


def _slide_favorability(prs, brand, period, sentiment, ai_report,
                         competitor_data=None, page_num=None):
    """好感度與負面聚焦：P/N 大數字 + 情緒占比條 + AI 負面議題 + 競品好感度對比"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, "好感度與負面聚焦", page_num, brand, period)

    pn = sentiment.get("pn_value")
    neg_pct = sentiment.get("negative_percent", "—")
    _numbered_insight(slide, ML, Inches(0.95), CW, Inches(0.55), [
        f"本期好感度（P/N 值）為 {pn}，負面聲量占比 {neg_pct}；"
        + ("正面聲量大於負面。" if isinstance(pn, (int, float)) and pn >= 1
           else "負面聲量已超越正面，建議加強正向議題操作。")
    ], size=12)

    # 左：P/N 卡 + 情緒占比 100% 條
    _data_card(slide, ML, Inches(1.7), Inches(5.9), Inches(1.15),
               str(pn), "好感度 P/N 值（正面 ÷ 負面）",
               accent=C_SLATE, value_size=34)

    bar_y = Inches(3.3)
    bar_w = Inches(5.9)
    pos = _safe_pct(sentiment.get("positive_percent"))
    neu = _safe_pct(sentiment.get("neutral_percent"))
    neg = _safe_pct(sentiment.get("negative_percent"))
    _tb(slide, ML, bar_y - Inches(0.32), bar_w, Inches(0.3),
        "情緒占比", size=12, bold=True, color=C_INK)
    x = ML
    for pct, color in [(pos, C_POS), (neu, C_NEU), (neg, C_NEG)]:
        seg = Emu(int(bar_w * (pct / 100))) if pct else Emu(0)
        if pct:
            _rect(slide, x, bar_y, seg, Inches(0.5), fill=color)
            if pct >= 8:
                _tb(slide, x, bar_y + Inches(0.1), seg, Inches(0.3),
                    f"{pct:.0f}%", size=10, bold=True, color=C_WHITE,
                    align=PP_ALIGN.CENTER, wrap=False)
            x = Emu(x + seg)
    leg_y = bar_y + Inches(0.62)
    for j, (lbl, color) in enumerate(
        [("正面", C_POS), ("中立", C_NEU), ("負面", C_NEG)]
    ):
        lx = ML + j * Inches(1.3)
        _rect(slide, lx, leg_y + Inches(0.03), Inches(0.18), Inches(0.18), fill=color)
        _tb(slide, lx + Inches(0.26), leg_y, Inches(1.0), Inches(0.28),
            lbl, size=10, color=C_GRAY, wrap=False)

    # 競品好感度對比（選填）
    if competitor_data:
        _tb(slide, ML, Inches(4.6), bar_w, Inches(0.3),
            "▎ 競品好感度對比", size=12, bold=True, color=C_INK)
        rows = [(brand, pn)] + [
            (c["name"], c["sentiment"].get("pn_value")) for c in competitor_data
        ]
        rows.sort(key=lambda r: r[1] if isinstance(r[1], (int, float)) else -1,
                  reverse=True)
        tbl = slide.shapes.add_table(
            len(rows) + 1, 2, ML, Inches(4.95), bar_w, Inches(0.4 * len(rows) + 0.4)
        ).table
        tbl.columns[0].width = Inches(3.9)
        tbl.columns[1].width = Inches(2.0)
        _cell(tbl.cell(0, 0), "品牌", size=11, bold=True, color=C_WHITE,
              bg=C_SLATE2)
        _cell(tbl.cell(0, 1), "P/N 值", size=11, bold=True, color=C_WHITE,
              bg=C_SLATE2, align=PP_ALIGN.CENTER)
        for i, (name, val) in enumerate(rows, start=1):
            bg = C_LIGHT if name == brand else (C_LIGHT if i % 2 == 0 else C_WHITE)
            _cell(tbl.cell(i, 0), name, bg=bg, bold=(name == brand))
            _cell(tbl.cell(i, 1), str(val), bg=bg, align=PP_ALIGN.CENTER)

    # 右：負面焦點（AI 負面議題）
    rx = Inches(7.1)
    _rect(slide, rx, Inches(1.7), Inches(5.7), Inches(0.5), fill=C_NEG)
    _tb(slide, rx, Inches(1.78), Inches(5.7), Inches(0.34),
        "負面焦點（AI 議題彙整）", size=13, bold=True, color=C_WHITE,
        align=PP_ALIGN.CENTER, wrap=False)

    secs = _parse_ai(ai_report)
    neg_lines = [
        _strip_md(l.strip().lstrip("-•0123456789.、 "))
        for l in secs.get("負面議題", "").split('\n') if l.strip()
    ][:5]
    if not neg_lines:
        neg_lines = ["本期未偵測到明顯負面議題，負面聲量維持低水位。"]
    _rect(slide, rx, Inches(2.2), Inches(5.7), Inches(4.4), fill=C_LIGHT)
    _bullet_box(slide, rx + Inches(0.2), Inches(2.4), Inches(5.3), Inches(4.0),
                [f"• {l[:60]}" for l in neg_lines], size=12, spacing=10, color=C_INK)


def _slide_data_summary(prs, brand, period, trend, growth_info=None,
                        competitor_data=None, page_num=None):
    """數據摘要：本期聲量 / 成長率 / 上期 大字呈現
    （對照官方月報結論章的『數據摘要』頁）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, "數據摘要", page_num, brand, period)

    days = len(trend["daily"]) or 1
    mom = (growth_info or {}).get("mom")
    yoy = (growth_info or {}).get("yoy")

    _numbered_insight(slide, ML, Inches(0.95), CW, Inches(0.55), [
        f"本期{brand}總聲量 {trend['total']:,} 筆，"
        f"日均 {trend['total'] // days:,} 筆。"
    ], size=12)

    # 三張大數字卡：本期聲量 / 本期成長率 / 年增率
    card_w = Inches(3.9)
    card_h = Inches(1.85)
    gap = Inches(0.31)
    cy = Inches(1.75)

    _data_card(slide, ML, cy, card_w, card_h,
               f"{trend['total']:,}", "本期總聲量（筆）",
               accent=C_SLATE, value_size=34)

    if mom:
        sign = "+" if mom["rate"] >= 0 else ""
        acc = C_POS if mom["rate"] >= 0 else C_NEG
        _data_card(slide, ML + card_w + gap, cy, card_w, card_h,
                   f"{sign}{mom['rate']}%", "本期聲量成長率（較上期）",
                   accent=acc, value_size=32)
    else:
        _data_card(slide, ML + card_w + gap, cy, card_w, card_h,
                   "—", "本期聲量成長率（無上期資料）",
                   accent=C_NEU, value_size=32)

    if yoy:
        sign = "+" if yoy["rate"] >= 0 else ""
        acc = C_POS if yoy["rate"] >= 0 else C_NEG
        _data_card(slide, ML + 2 * (card_w + gap), cy, card_w, card_h,
                   f"{sign}{yoy['rate']}%", "年增率（較去年同期）",
                   accent=acc, value_size=32)
    else:
        _data_card(slide, ML + 2 * (card_w + gap), cy, card_w, card_h,
                   "—", "年增率（無去年同期資料）",
                   accent=C_NEU, value_size=32)

    # 對照基準
    refs = []
    if mom:
        refs.append(f"上期聲量：{mom['prev_total']:,} 筆")
    if yoy:
        refs.append(f"去年同期聲量：{yoy['prev_total']:,} 筆")
    if refs:
        _tb(slide, ML, cy + card_h + Inches(0.18), CW, Inches(0.4),
            "　｜　".join(refs), size=13, color=C_GRAY)

    # 競品本期聲量對比（選填）
    if competitor_data:
        _tb(slide, ML, Inches(4.55), CW, Inches(0.3),
            "▎ 本期聲量排名（含競品）", size=12, bold=True, color=C_INK)
        ranked = sorted(
            [{"name": brand, "total": trend["total"], "self": True}]
            + [{"name": c["name"], "total": c["trend"]["total"], "self": False}
               for c in competitor_data],
            key=lambda x: x["total"], reverse=True,
        )
        tbl = slide.shapes.add_table(
            len(ranked) + 1, 3, ML, Inches(4.9), CW,
            Inches(0.4 * len(ranked) + 0.4)
        ).table
        for i, cw in enumerate([Inches(1.5), Inches(7.5), Inches(3.33)]):
            tbl.columns[i].width = cw
        for j, hdr in enumerate(["排名", "品牌", "本期總聲量（筆）"]):
            _cell(tbl.cell(0, j), hdr, size=11, bold=True, color=C_WHITE,
                  bg=C_SLATE2, align=PP_ALIGN.CENTER)
        for i, r in enumerate(ranked, start=1):
            bg = C_LIGHT if r["self"] else (C_LIGHT if i % 2 == 0 else C_WHITE)
            _cell(tbl.cell(i, 0), str(i), bg=bg, align=PP_ALIGN.CENTER)
            _cell(tbl.cell(i, 1), r["name"], bg=bg, bold=r["self"])
            _cell(tbl.cell(i, 2), f"{r['total']:,}", bg=bg, align=PP_ALIGN.CENTER)


def _slide_observation_matrix(prs, brand, period, trend, sentiment, keywords,
                              kols, growth_info=None, page_num=None):
    """觀測彙整：以固定觀測面向（聲量／情緒／話題／傳播／負面）條列本期結論
    （對照官方月報結論章的『觀測彙整』）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, "本期觀測彙整", page_num, brand, period)

    days = len(trend["daily"]) or 1
    avg = trend["total"] // days

    growth_txt = ""
    if growth_info:
        mom = growth_info.get("mom")
        if mom:
            word = "成長" if mom["rate"] >= 0 else "下滑"
            growth_txt = (f"，較上期（{mom['prev_total']:,}筆）{word} "
                          f"{abs(mom['rate'])}%")

    top_kw = "、".join(
        f"{k['keyword']}（{k['count']}）" for k in keywords[:3]
    ) if keywords else "—"
    top_kol = (f"{kols[0]['platform']} {kols[0]['channel']}"
               f"（互動 {kols[0].get('engagement', 0):,}）") if kols else "—"

    rows = [
        ("聲量表現", C_SLATE,
         f"本期總聲量 {trend['total']:,}筆，日均 {avg:,}筆{growth_txt}；"
         f"聲量高峰 {trend['peak_date']}（{trend['peak_volume']:,}筆）。"),
        ("情緒結構", C_BLUE,
         f"正面 {sentiment['positive_percent']}、中立 {sentiment['neutral_percent']}、"
         f"負面 {sentiment['negative_percent']}；好感度 P/N 值 {sentiment['pn_value']}。"),
        ("熱門話題", C_SLATE2,
         f"本期關鍵字 TOP3：{top_kw}。"),
        ("關鍵傳播", C_SLATE,
         f"最具影響力 KOL：{top_kol}，為本期聲量主要擴散來源。"),
        ("負面聚焦", C_NEG,
         f"負面聲量 {sentiment['negative']:,}筆（{sentiment['negative_percent']}），"
         + ("超過 100 筆，建議持續監控社群討論。"
            if sentiment["negative"] > 100 else "維持低水位。")),
    ]

    y = Inches(1.1)
    row_h = Inches(1.05)
    label_w = Inches(1.9)
    for i, (label, color, text) in enumerate(rows):
        ry = y + i * (row_h + Inches(0.08))
        _rect(slide, ML, ry, label_w, row_h, fill=color)
        _tb(slide, ML, ry + Inches(0.34), label_w, Inches(0.4),
            label, size=14, bold=True, color=C_WHITE,
            align=PP_ALIGN.CENTER, wrap=False)
        _rect(slide, ML + label_w + Inches(0.12), ry,
              CW - label_w - Inches(0.12), row_h, fill=C_LIGHT)
        _tb(slide, ML + label_w + Inches(0.35), ry + Inches(0.16),
            CW - label_w - Inches(0.6), row_h - Inches(0.3),
            text, size=13, color=C_INK)


def _slide_keyword_events(prs, brand, period, keyword_events, page_num):
    """關鍵字對應事件：把熱門關鍵字還原成背後的代表性文章／事件
    （對照官方月報『熱門關鍵字…對應事件如下』）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, "關鍵字對應事件", page_num, brand, period)

    _numbered_insight(slide, ML, Inches(0.95), CW, Inches(0.55), [
        "下表將本期熱門關鍵字對應到互動最高的代表性文章，"
        "協助把關鍵字還原成具體事件脈絡。"
    ], size=12)

    rows = [ke for ke in keyword_events if ke.get("keyword")][:8]
    if not rows:
        return

    n = len(rows)
    tbl = slide.shapes.add_table(
        n + 1, 4, ML, Inches(1.65), CW, Inches(0.6 * n + 0.5)
    ).table
    for i, cw in enumerate(
        [Inches(2.0), Inches(1.2), Inches(1.4), Inches(7.73)]
    ):
        tbl.columns[i].width = cw

    for j, hdr in enumerate(["關鍵字", "聲量", "命中文章", "對應熱門文章／事件"]):
        _cell(tbl.cell(0, j), hdr, size=11, bold=True, color=C_WHITE,
              bg=C_SLATE2, align=PP_ALIGN.CENTER)

    for i, ke in enumerate(rows, start=1):
        bg = C_LIGHT if i % 2 == 0 else C_WHITE
        art = ke.get("article")
        if art:
            event = f"[{art['source']}] {_clean(art.get('title', ''), 40)}"
        else:
            event = "（本期熱門文章未見直接對應）"
        _cell(tbl.cell(i, 0), ke["keyword"], bg=bg, bold=True)
        _cell(tbl.cell(i, 1), f"{ke.get('count', 0):,}", bg=bg,
              align=PP_ALIGN.CENTER)
        _cell(tbl.cell(i, 2), f"{ke.get('match_count', 0)} 篇", bg=bg,
              align=PP_ALIGN.CENTER)
        _cell(tbl.cell(i, 3), event, bg=bg)


def _slide_keywords(prs, brand, period, keywords, charts_dir, page_num=5, caption=None):
    """熱門關鍵字橫條圖"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, "熱門關鍵字 TOP10", page_num, brand, period)

    insight_items = []
    if keywords:
        top = keywords[0]
        insight_items.append(
            f"本期最熱關鍵字為「{top['keyword']}」，累計聲量 {top['count']:,} 筆，"
            f"建議對照熱門文章確認討論脈絡。"
        )
    if caption:
        insight_items.append(_clean(caption, 130))
    _numbered_insight(slide, ML, Inches(0.95), CW, Inches(1.3), insight_items, size=12)

    slide.shapes.add_picture(
        f"{charts_dir}/keyword_bar.png",
        ML, Inches(2.4), width=CW, height=Inches(4.5))


def _slide_wordcloud(prs, brand, period, charts_dir, page_num=6):
    """關鍵字雲"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, "關鍵字雲", page_num, brand, period)

    slide.shapes.add_picture(
        f"{charts_dir}/keyword_cloud.png",
        ML, CY, width=CW, height=Inches(5.9))


def _slide_articles(prs, brand, period, articles, page_num=7, caption=None):
    """熱門文章分析：圈碼洞察 + TOP3 話題卡片 + TOP4-10 精簡表格"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, "熱門文章分析", page_num, brand, period)

    top3 = articles[:3]
    rest = articles[3:10]

    insight_items = []
    if top3:
        insight_items.append(
            f"本期互動數最高文章來自 [{top3[0]['source']}] "
            f"{_clean(top3[0]['channel'], 14)}（互動 {top3[0]['engagement']:,}）。"
        )
    if caption:
        insight_items.append(_clean(caption, 130))
    if insight_items:
        _numbered_insight(slide, ML, Inches(0.95), CW, Inches(1.2), insight_items, size=12)

    top_y = Inches(2.25)

    # ── TOP3 話題卡片 ────────────────────────────────
    gap = Inches(0.2)
    card_w = Emu(int((CW - gap * 2) / 3))
    card_h = Inches(1.95)

    for i, art in enumerate(top3):
        cx = ML + i * (card_w + gap)
        cls = classify_article(art, brand)
        _rect(slide, cx, top_y, card_w, card_h, fill=C_LIGHT)
        _rect(slide, cx, top_y, card_w, Inches(0.06), fill=RANK_COLORS[i])
        _rank_badge(slide, cx + Inches(0.4), top_y + Inches(0.42),
                    Inches(0.5), i + 1, size=15)
        _tb(slide, cx + Inches(0.75), top_y + Inches(0.2),
            card_w - Inches(1.85), Inches(0.5),
            f"{art['source']} ・ {_clean(art['channel'], 8)}",
            size=10, bold=True, color=C_SLATE, wrap=True)
        # 分類標籤（右上角色塊）
        pill_w = Inches(0.92)
        _rect(slide, cx + card_w - pill_w - Inches(0.14), top_y + Inches(0.18),
              pill_w, Inches(0.28),
              fill=CATEGORY_COLORS.get(cls, C_SLATE2))
        _tb(slide, cx + card_w - pill_w - Inches(0.14), top_y + Inches(0.2),
            pill_w, Inches(0.24), cls, size=8.5, bold=True, color=C_WHITE,
            align=PP_ALIGN.CENTER, wrap=False)
        _tb(slide, cx + Inches(0.18), top_y + Inches(0.8),
            card_w - Inches(0.36), Inches(0.75),
            _clean(art["title"], 46), size=12, color=C_INK)
        _tb(slide, cx + Inches(0.18), top_y + card_h - Inches(0.36),
            card_w - Inches(0.36), Inches(0.32),
            f"互動 {art['engagement']:,}　讚 {art['likes']:,}　分享 {art['replies']:,}",
            size=9.5, color=C_GRAY, wrap=False)

    # ── TOP4–10 精簡表格 ─────────────────────────────
    if rest:
        table_y = top_y + card_h + Inches(0.18)
        n = len(rest)
        tbl = slide.shapes.add_table(
            n + 1, 6, ML, table_y, CW, Inches(2.55)
        ).table

        for i, cw in enumerate(
            [Inches(0.55), Inches(1.1), Inches(1.4), Inches(1.8),
             Inches(5.7), Inches(1.78)]
        ):
            tbl.columns[i].width = cw

        for j, hdr in enumerate(
            ["排名", "平台", "分類", "頻道", "標題", "互動（cc）"]
        ):
            _cell(tbl.cell(0, j), hdr,
                  size=11, bold=True, color=C_WHITE,
                  bg=C_SLATE2, align=PP_ALIGN.CENTER)

        for i, art in enumerate(rest, start=1):
            bg = C_LIGHT if i % 2 == 0 else C_WHITE
            cls = classify_article(art, brand)
            _cell(tbl.cell(i, 0), str(art["rank"]),
                  bg=bg, align=PP_ALIGN.CENTER)
            _cell(tbl.cell(i, 1), art["source"],
                  bg=bg, align=PP_ALIGN.CENTER)
            _cell(tbl.cell(i, 2), cls,
                  bg=bg, align=PP_ALIGN.CENTER)
            _cell(tbl.cell(i, 3), _clean(art["channel"], 12),
                  bg=bg)
            _cell(tbl.cell(i, 4), _clean(art["title"], 40),
                  bg=bg)
            _cell(tbl.cell(i, 5), f"{art['engagement']:,}",
                  bg=bg, align=PP_ALIGN.CENTER)


def _slide_news(prs, brand, period, articles, page_num, caption=None):
    """新聞文章列表：篩出分類為「新聞媒體」的文章，依互動排序列表
    （對照 KEYPO 官方月報的『新聞文章列表』頁）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, "新聞文章列表", page_num, brand, period)

    news = [a for a in articles if classify_article(a, brand) == "新聞媒體"]
    news = sorted(news, key=lambda a: a.get("engagement", 0), reverse=True)[:8]

    insight_items = []
    if news:
        insight_items.append(
            f"本期新聞媒體共 {len([a for a in articles if classify_article(a, brand) == '新聞媒體'])} 篇進入熱門討論，"
            f"以下依互動數列出前 {len(news)} 篇。"
        )
    else:
        insight_items.append("本期熱門文章中未偵測到新聞媒體來源，新聞聲量較低。")
    if caption:
        insight_items.append(_clean(caption, 120))
    _numbered_insight(slide, ML, Inches(0.95), CW, Inches(1.1), insight_items, size=12)

    if not news:
        return

    n = len(news)
    tbl = slide.shapes.add_table(
        n + 1, 4, ML, Inches(2.2), CW, Inches(0.5 * n + 0.5)
    ).table
    for i, cw in enumerate(
        [Inches(1.25), Inches(1.75), Inches(7.55), Inches(1.78)]
    ):
        tbl.columns[i].width = cw

    for j, hdr in enumerate(["日期", "來源", "標題", "互動數"]):
        _cell(tbl.cell(0, j), hdr, size=11, bold=True, color=C_WHITE,
              bg=C_SLATE2, align=PP_ALIGN.CENTER)

    for i, art in enumerate(news, start=1):
        bg = C_LIGHT if i % 2 == 0 else C_WHITE
        date = str(art.get("time", "")).split(" ")[0]
        _cell(tbl.cell(i, 0), date, bg=bg, align=PP_ALIGN.CENTER)
        _cell(tbl.cell(i, 1),
              f"{art['source']} / {_clean(art['channel'], 8)}", bg=bg)
        _cell(tbl.cell(i, 2), _clean(art["title"], 46), bg=bg)
        _cell(tbl.cell(i, 3), f"{art['engagement']:,}",
              bg=bg, align=PP_ALIGN.CENTER)


def _slide_kol(prs, brand, period, kols, page_num=8, caption=None):
    """網路關鍵領袖 TOP10 四維指標表（發表數／回文／按讚／互動總計，
    排名欄以徽章色塊標示前三名）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, "網路關鍵領袖 TOP10", page_num, brand, period)

    insight_items = []
    if kols:
        top = kols[0]
        insight_items.append(
            f"本期最具影響力 KOL 為 [{top['platform']}] "
            f"{_clean(top['channel'], 18)}（互動總計 {top.get('engagement', 0):,}、"
            f"按讚 {top.get('likes', 0):,}）。"
        )
    if caption:
        insight_items.append(_clean(caption, 130))
    if insight_items:
        _numbered_insight(slide, ML, Inches(0.95), CW, Inches(1.2), insight_items, size=12)

    n = min(len(kols), 10)
    tbl = slide.shapes.add_table(
        n + 1, 7,
        ML, Inches(2.25),
        CW, Inches(4.65)
    ).table

    # 欄寬（合計 = 12.33"）
    for i, cw in enumerate(
        [Inches(0.6), Inches(1.4), Inches(4.28), Inches(1.1),
         Inches(1.35), Inches(1.5), Inches(2.1)]
    ):
        tbl.columns[i].width = cw

    for j, hdr in enumerate(
        ["排名", "平台", "頻道名稱", "發表數", "回文數", "按讚數", "互動總計"]
    ):
        _cell(tbl.cell(0, j), hdr,
              size=11, bold=True, color=C_WHITE,
              bg=C_SLATE2, align=PP_ALIGN.CENTER)

    for i, kol in enumerate(kols[:n], start=1):
        bg = C_LIGHT if i % 2 == 0 else C_WHITE
        rank_bg = RANK_COLORS[i - 1] if i <= 3 else bg
        rank_color = C_WHITE if i <= 3 else C_INK
        _cell(tbl.cell(i, 0), str(kol["rank"]),
              bg=rank_bg, bold=(i <= 3), color=rank_color,
              align=PP_ALIGN.CENTER)
        _cell(tbl.cell(i, 1), kol["platform"],
              bg=bg, align=PP_ALIGN.CENTER)
        _cell(tbl.cell(i, 2), _clean(kol["channel"], 24), bg=bg)
        _cell(tbl.cell(i, 3), f"{kol.get('posts', 0):,}",
              bg=bg, align=PP_ALIGN.CENTER)
        _cell(tbl.cell(i, 4), f"{kol.get('replies', 0):,}",
              bg=bg, align=PP_ALIGN.CENTER)
        _cell(tbl.cell(i, 5), f"{kol.get('likes', 0):,}",
              bg=bg, align=PP_ALIGN.CENTER)
        _cell(tbl.cell(i, 6), f"{kol.get('engagement', 0):,}",
              bg=bg, bold=True, align=PP_ALIGN.CENTER)


def _slide_competitor_volume(prs, brand, period, trend, competitor_data, charts_dir,
                              page_num, captions=None):
    """競品聲量比較：圈碼洞察 + 折線疊圖 + 各品牌事件摘要 + 排名小表格"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, "競品聲量分析", page_num, brand, period)

    ranked = sorted(
        [{"name": brand, "total": trend["total"]}]
        + [{"name": c["name"], "total": c["trend"]["total"]} for c in competitor_data],
        key=lambda x: x["total"], reverse=True
    )
    top = ranked[0]
    insight = f"本期聲量最高為「{top['name']}」（{top['total']:,}筆）。"
    if ranked[0]["name"] != brand:
        main_rank = next(i for i, r in enumerate(ranked, 1) if r["name"] == brand)
        insight += f"{brand} 本期排名第 {main_rank}。"
    _numbered_insight(slide, ML, Inches(0.95), CW, Inches(0.5), [insight], size=12)

    slide.shapes.add_picture(
        f"{charts_dir}/competitor_volume.png",
        ML, Inches(1.55), width=CW, height=Inches(2.2))

    if captions:
        caption_lines = []
        for c in competitor_data:
            text = captions.get(c["name"])
            if text:
                caption_lines.append(f"● {_clean(text, 110)}")
        if caption_lines:
            _bullet_box(slide, ML, Inches(3.85), CW, Inches(1.3),
                        caption_lines, size=9.5, spacing=3, color=C_GRAY)

    table_y = Inches(5.25)
    n = len(ranked)
    tbl = slide.shapes.add_table(n + 1, 3, ML, table_y, CW, Inches(1.65)).table
    for i, cw in enumerate([Inches(1.5), Inches(7.5), Inches(3.33)]):
        tbl.columns[i].width = cw
    for j, hdr in enumerate(["排名", "品牌", "總聲量（筆）"]):
        _cell(tbl.cell(0, j), hdr, size=11, bold=True, color=C_WHITE,
              bg=C_SLATE2, align=PP_ALIGN.CENTER)
    for i, r in enumerate(ranked, start=1):
        bg = C_LIGHT if r["name"] == brand else (C_LIGHT if i % 2 == 0 else C_WHITE)
        _cell(tbl.cell(i, 0), str(i), bg=bg, align=PP_ALIGN.CENTER)
        _cell(tbl.cell(i, 1), r["name"], bg=bg, bold=(r["name"] == brand))
        _cell(tbl.cell(i, 2), f"{r['total']:,}", bg=bg, align=PP_ALIGN.CENTER)


def _slide_competitor_sentiment(prs, brand, period, sentiment, competitor_data, charts_dir, page_num):
    """競品情緒分布比較：圈碼洞察 + 100% 堆疊橫條圖"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, "競品情緒分布比較", page_num, brand, period)

    pn_values = [(brand, sentiment["pn_value"])] + [
        (c["name"], c["sentiment"]["pn_value"]) for c in competitor_data
    ]
    best = max(pn_values, key=lambda x: x[1] if isinstance(x[1], (int, float)) else 0)
    _numbered_insight(slide, ML, Inches(0.95), CW, Inches(0.5), [
        f"本期 P/N 值最高為「{best[0]}」（{best[1]}），"
        f"{brand} 本期 P/N 值為 {sentiment['pn_value']}。"
    ], size=12)

    slide.shapes.add_picture(
        f"{charts_dir}/competitor_sentiment.png",
        ML, Inches(1.55), width=CW, height=Inches(5.2))


def _slide_competitor_articles(prs, brand, period, competitor_data, page_num):
    """競品熱門文章：各競品 TOP3 合併表格"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, "競品熱門文章", page_num, brand, period)

    _numbered_insight(slide, ML, Inches(0.95), CW, Inches(0.5), [
        "以下列出各競品本期互動數最高的 TOP3 文章，供對照本品牌熱門文章分析。",
    ], size=12)

    rows = []
    for comp in competitor_data:
        for art in comp["articles"][:3]:
            rows.append((comp["name"], art))

    if not rows:
        return

    tbl = slide.shapes.add_table(
        len(rows) + 1, 6, ML, Inches(1.55), CW, Inches(5.3)
    ).table
    for i, cw in enumerate(
        [Inches(1.5), Inches(1.1), Inches(1.3), Inches(1.7),
         Inches(4.9), Inches(1.83)]
    ):
        tbl.columns[i].width = cw
    for j, hdr in enumerate(["品牌", "平台", "分類", "頻道", "標題", "互動（cc）"]):
        _cell(tbl.cell(0, j), hdr, size=11, bold=True, color=C_WHITE,
              bg=C_SLATE2, align=PP_ALIGN.CENTER)
    for i, (name, art) in enumerate(rows, start=1):
        bg = C_LIGHT if i % 2 == 0 else C_WHITE
        _cell(tbl.cell(i, 0), name, bg=bg, bold=True)
        _cell(tbl.cell(i, 1), art["source"], bg=bg, align=PP_ALIGN.CENTER)
        _cell(tbl.cell(i, 2), classify_article(art, name), bg=bg,
              align=PP_ALIGN.CENTER)
        _cell(tbl.cell(i, 3), _clean(art["channel"], 12), bg=bg)
        _cell(tbl.cell(i, 4), _clean(art["title"], 38), bg=bg)
        _cell(tbl.cell(i, 5), f"{art['engagement']:,}", bg=bg, align=PP_ALIGN.CENTER)


def _slide_competitor_kol(prs, brand, period, competitor_data, page_num):
    """競品關鍵領袖：各競品 TOP3 KOL 四維指標合併表
    （由各競品熱門文章依發文者聚合，對照官方月報各競品的關鍵領袖 TOP5）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, "競品關鍵領袖", page_num, brand, period)

    _numbered_insight(slide, ML, Inches(0.95), CW, Inches(0.5), [
        "以下彙整各競品本期互動最高的關鍵領袖（依熱門文章發文者聚合），"
        "供對照本品牌 KOL 分析。",
    ], size=12)

    rows = []
    for comp in competitor_data:
        for kol in parse_kol_from_articles(comp["articles"], limit=3):
            rows.append((comp["name"], kol))

    if not rows:
        return

    n = len(rows)
    tbl = slide.shapes.add_table(
        n + 1, 7, ML, Inches(1.55), CW, Inches(0.48 * n + 0.5)
    ).table
    for i, cw in enumerate(
        [Inches(1.5), Inches(1.2), Inches(3.7), Inches(1.0),
         Inches(1.3), Inches(1.5), Inches(2.13)]
    ):
        tbl.columns[i].width = cw
    for j, hdr in enumerate(
        ["品牌", "平台", "頻道名稱", "發表數", "回文數", "按讚數", "互動總計"]
    ):
        _cell(tbl.cell(0, j), hdr, size=11, bold=True, color=C_WHITE,
              bg=C_SLATE2, align=PP_ALIGN.CENTER)
    for i, (name, kol) in enumerate(rows, start=1):
        bg = C_LIGHT if i % 2 == 0 else C_WHITE
        _cell(tbl.cell(i, 0), name, bg=bg, bold=True)
        _cell(tbl.cell(i, 1), kol["platform"], bg=bg, align=PP_ALIGN.CENTER)
        _cell(tbl.cell(i, 2), _clean(kol["channel"], 22), bg=bg)
        _cell(tbl.cell(i, 3), f"{kol.get('posts', 0):,}", bg=bg,
              align=PP_ALIGN.CENTER)
        _cell(tbl.cell(i, 4), f"{kol.get('replies', 0):,}", bg=bg,
              align=PP_ALIGN.CENTER)
        _cell(tbl.cell(i, 5), f"{kol.get('likes', 0):,}", bg=bg,
              align=PP_ALIGN.CENTER)
        _cell(tbl.cell(i, 6), f"{kol.get('engagement', 0):,}", bg=bg,
              bold=True, align=PP_ALIGN.CENTER)


def _slide_ai(prs, brand, period, ai_report, page_num=9):
    """AI 輿情洞察（雙欄版面）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, "AI 輿情洞察", page_num, brand, period)

    secs = _parse_ai(ai_report)

    col_w = Inches(5.9)
    gap   = Inches(0.43)
    lx    = ML
    rx    = ML + col_w + gap

    # ── 左欄 ─────────────────────────────────────────
    # 執行摘要
    _section_label(slide, lx, Inches(0.95), "執行摘要")
    _rect(slide, lx, Inches(1.28), col_w, Inches(1.62), fill=C_LIGHT)
    summary = _strip_md(secs.get("執行摘要", ""))[:230]
    _bullet_box(slide, lx + Inches(0.15), Inches(1.35),
                col_w - Inches(0.2), Inches(1.5),
                [summary], size=11)

    # 重點洞察
    _section_label(slide, lx, Inches(3.02), "重點洞察")
    raw_lines = [
        _strip_md(l.strip())
        for l in secs.get("重點洞察", "").split('\n')
        if l.strip() and not l.strip().startswith('(')
    ]
    _numbered_insight(slide, lx, Inches(3.38),
                       col_w, Inches(3.52),
                       [l[:78] + ('…' if len(l) > 78 else '') for l in raw_lines[:4]],
                       size=12, spacing=8)

    # ── 右欄 ─────────────────────────────────────────
    # 正面議題
    _section_label(slide, rx, Inches(0.95), "正面議題", color=C_POS)
    pos_lines = [
        _strip_md(l.strip())
        for l in secs.get("正面議題", "").split('\n')
        if l.strip()
    ][:3]
    _bullet_box(slide, rx, Inches(1.28),
                col_w, Inches(1.55),
                [f"● {l[:70]}" for l in pos_lines],
                size=12, spacing=6)

    _rect(slide, rx, Inches(2.9), col_w, Inches(0.02), fill=C_LINE)

    # 負面議題
    _section_label(slide, rx, Inches(2.97), "負面議題", color=C_NEG)
    neg_lines = [
        _strip_md(l.strip())
        for l in secs.get("負面議題", "").split('\n')
        if l.strip()
    ][:3]
    _bullet_box(slide, rx, Inches(3.32),
                col_w, Inches(1.42),
                [f"● {l[:70]}" for l in neg_lines],
                size=12, spacing=6)

    _rect(slide, rx, Inches(4.8), col_w, Inches(0.02), fill=C_LINE)

    # 建議行動
    _section_label(slide, rx, Inches(4.87), "建議行動", color=C_BLUE)
    act_lines = [
        _strip_md(l.strip())
        for l in secs.get("建議行動", "").split('\n')
        if l.strip()
    ][:4]
    _bullet_box(slide, rx, Inches(5.22),
                col_w, Inches(1.7),
                [f"● {l[:70]}" for l in act_lines],
                size=12, spacing=6)


def _slide_conclusion(prs, brand, period, ai_report, page_num=10):
    """結論：3 欄卡片（重點洞察／風險觀察／建議行動）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, "結論與建議", page_num, brand, period)

    secs = _parse_ai(ai_report)
    columns = [
        ("重點洞察", secs.get("重點洞察", "")),
        ("風險觀察", secs.get("風險觀察", "")),
        ("建議行動", secs.get("建議行動", "")),
    ]

    gap = Inches(0.25)
    col_w = Emu(int((CW - gap * 2) / 3))
    header_h = Inches(0.6)
    body_h = Inches(4.6)
    top_y = Inches(1.2)

    for i, (title, text) in enumerate(columns):
        cx = ML + i * (col_w + gap)
        _rect(slide, cx, top_y, col_w, header_h, fill=C_BLUE)
        _tb(slide, cx, top_y + Inches(0.12), col_w, Inches(0.4),
            title, size=16, bold=True, color=C_WHITE,
            align=PP_ALIGN.CENTER, wrap=False)
        _rect(slide, cx, top_y + header_h, col_w, body_h, fill=C_BLUE_L)

        lines = [
            _strip_md(l.strip().lstrip("-•0123456789.、 "))
            for l in text.split('\n') if l.strip()
        ][:5]
        if not lines:
            lines = ["（無資料）"]

        _bullet_box(slide, cx + Inches(0.2), top_y + header_h + Inches(0.2),
                    col_w - Inches(0.4), body_h - Inches(0.4),
                    [f"• {l[:60]}" for l in lines],
                    size=12, spacing=8, color=C_INK)


def _slide_back(prs, brand, period):
    """封底：白底 + 對角裝飾（與封面呼應）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _diagonal(slide, -Inches(1.8), -Inches(3.8), Inches(6.0), C_LIGHT)
    _diagonal(slide, W - Inches(1.4), H - Inches(1.4), Inches(4.0), C_RED)
    _diagonal(slide, -Inches(2.4), H - Inches(1.6), Inches(3.0), C_RED)

    _tb(slide, Inches(0.9), Inches(2.5), Inches(11.0), Inches(1.1),
        brand, size=38, bold=True, color=C_INK)
    _tb(slide, Inches(0.9), Inches(3.7), Inches(10.0), Inches(0.6),
        "網路輿情分析報告", size=20, color=C_RED)
    _rect(slide, Inches(0.9), Inches(4.42), Inches(8.5), Inches(0.03), fill=C_LINE)
    _tb(slide, Inches(0.9), Inches(4.57), Inches(10.0), Inches(0.4),
        period, size=13, color=C_GRAY)
    _tb(slide, Inches(0.9), Inches(5.08), Inches(10.0), Inches(0.35),
        "資料來源｜KEYPO 大數據關鍵引擎", size=11, color=C_GRAY)
    _tb(slide, Inches(0.9), Inches(6.45), Inches(10.0), Inches(0.35),
        "— 報告完 —", size=12, color=C_SLATE)


# ══════════════════════════════════════════════════════
# 主函式
# ══════════════════════════════════════════════════════

def generate_ppt(
    brand_name, report, ai_report,
    keywords, kols, articles,
    sentiment, trend,
    start_date, end_date,
    out_dir="output",
    growth_info=None,
    competitor_data=None,
    table_captions=None,
    source_dist=None,
    trend_narrative=None
):
    charts_dir = f"{out_dir}/charts"
    os.makedirs(out_dir, exist_ok=True)
    competitor_data = competitor_data or []
    table_captions = table_captions or {}

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    disp_s = format_date(start_date)
    disp_e = format_date(end_date)
    period = f"{disp_s} – {disp_e}"

    # page：連續頁碼（封面/封底不編號）；chapter：分節頁章節序號，兩者獨立。
    page = 1
    chapter = 0
    _slide_cover(prs, brand_name, period)

    # ── 開場：調查方法 / 名詞說明 ─────────────────────
    page += 1
    _slide_method(prs, brand_name, period, page)
    page += 1
    _slide_glossary(prs, brand_name, period, page)

    # ── CHAPTER 1 網路輿情總覽 ────────────────────────
    chapter += 1
    page += 1
    _section_divider(prs, "網路輿情總覽",
                     f"SOCIAL LISTENING OVERVIEW　｜　{period}", chapter, page)
    page += 1
    _slide_summary(prs, brand_name, period, trend, sentiment, keywords,
                   growth_info=growth_info, page_num=page)
    page += 1
    _slide_trend(prs, brand_name, period, trend, articles, charts_dir,
                 growth_info=growth_info, page_num=page,
                 narrative=trend_narrative)
    if source_dist:
        page += 1
        _slide_source(prs, brand_name, period, source_dist, charts_dir, page)
    page += 1
    _slide_sentiment(prs, brand_name, period, sentiment, charts_dir, page_num=page)
    page += 1
    _slide_favorability(prs, brand_name, period, sentiment, ai_report,
                        competitor_data=competitor_data, page_num=page)
    page += 1
    _slide_keywords(prs, brand_name, period, keywords, charts_dir, page_num=page,
                     caption=table_captions.get("keywords"))
    page += 1
    _slide_keyword_events(prs, brand_name, period,
                          match_keywords_to_articles(keywords, articles), page)
    page += 1
    _slide_wordcloud(prs, brand_name, period, charts_dir, page_num=page)
    page += 1
    _slide_articles(prs, brand_name, period, articles, page_num=page,
                     caption=table_captions.get("articles"))
    page += 1
    _slide_news(prs, brand_name, period, articles, page_num=page,
                caption=table_captions.get("articles"))
    page += 1
    _slide_kol(prs, brand_name, period, kols, page_num=page,
               caption=table_captions.get("kol"))

    # ── CHAPTER 2 競品聲量分析（選填）─────────────────
    if competitor_data:
        chapter += 1
        page += 1
        _section_divider(prs, "競品聲量分析",
                          f"COMPETITOR ANALYSIS　｜　{period}", chapter, page)
        page += 1
        _slide_competitor_volume(prs, brand_name, period, trend, competitor_data,
                                  charts_dir, page_num=page,
                                  captions=table_captions.get("competitors"))
        page += 1
        _slide_competitor_sentiment(prs, brand_name, period, sentiment, competitor_data,
                                     charts_dir, page_num=page)
        page += 1
        _slide_competitor_articles(prs, brand_name, period, competitor_data,
                                    page_num=page)
        page += 1
        _slide_competitor_kol(prs, brand_name, period, competitor_data,
                              page_num=page)

    # ── CHAPTER 3 結論 ────────────────────────────────
    chapter += 1
    page += 1
    _section_divider(prs, "結論與建議",
                     f"CONCLUSION　｜　{period}", chapter, page)
    page += 1
    _slide_data_summary(prs, brand_name, period, trend, growth_info=growth_info,
                        competitor_data=competitor_data, page_num=page)
    page += 1
    _slide_observation_matrix(prs, brand_name, period, trend, sentiment, keywords,
                              kols, growth_info=growth_info, page_num=page)
    page += 1
    _slide_ai(prs, brand_name, period, ai_report, page_num=page)
    page += 1
    _slide_conclusion(prs, brand_name, period, ai_report, page_num=page)
    _slide_back(prs, brand_name, period)

    timestamp = datetime.now().strftime("%H%M%S")
    output_path = (
        f"{out_dir}/{brand_name}_{start_date}_{end_date}"
        f"_{timestamp}_Report.pptx"
    )
    prs.save(output_path)

    print(f"PPT 已儲存：{output_path}")
    if os.path.exists(output_path):
        os.startfile(os.path.abspath(output_path))

    return output_path
